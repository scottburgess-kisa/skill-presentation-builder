#!/usr/bin/env python3
"""
Create new DEFRA PowerPoint presentations from scratch based on user content.
"""

import os
import sys
from datetime import datetime
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor

def create_defra_presentation():
    """Create a new presentation with DEFRA template."""
    prs = Presentation()
    prs.slide_width = Cm(25.4)
    prs.slide_height = Cm(14.29)
    return prs

def add_background_image(slide, image_path):
    """Add background image to slide."""
    try:
        slide.shapes.add_picture(image_path, 0, 0, Cm(25.4), Cm(14.29))
    except:
        # If image not found, add colored background
        background = slide.shapes.add_shape(1, 0, 0, Cm(25.4), Cm(14.29))
        background.fill.solid()
        background.fill.fore_color.rgb = RGBColor(0, 175, 65)  # DEFRA Green
        background.line.fill.background()

def add_logo(slide, logo_path, x, y, w, h):
    """Add DEFRA logo to slide."""
    try:
        slide.shapes.add_picture(logo_path, Cm(x), Cm(y), Cm(w), Cm(h))
    except:
        # If logo not found, add placeholder
        logo_placeholder = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
        logo_placeholder.fill.solid()
        logo_placeholder.fill.fore_color.rgb = RGBColor(235, 235, 235)
        logo_placeholder.line.color.rgb = RGBColor(0, 175, 65)

def create_cover_slide(prs, title_text, project_name="", version="1.0"):
    """Create DEFRA cover slide."""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Background
    add_background_image(slide, "assets/cover_bg.png")

    # DEFRA logo
    add_logo(slide, "assets/logo_wide.png", 19.29, 0.38, 5.65, 1.50)

    # Title
    title_box = slide.shapes.add_textbox(Cm(1.20), Cm(4.52), Cm(16.89), Cm(2.65))
    title_frame = title_box.text_frame
    title_frame.text = title_text

    # Format all paragraphs in title as white text
    for paragraph in title_frame.paragraphs:
        paragraph.font.name = "Arial"
        paragraph.font.size = Pt(28)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Supplier
    supplier_box = slide.shapes.add_textbox(Cm(1.20), Cm(7.20), Cm(19.98), Cm(1.03))
    supplier_frame = supplier_box.text_frame
    supplier_frame.text = "Equal Experts"
    supplier_frame.paragraphs[0].font.name = "Arial"
    supplier_frame.paragraphs[0].font.size = Pt(24)
    supplier_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Date/Version info
    info_box = slide.shapes.add_textbox(Cm(1.20), Cm(8.83), Cm(19.98), Cm(1.50))
    info_frame = info_box.text_frame
    today = datetime.now().strftime("%d %b %Y")
    info_frame.text = f"Date       {today}         Version    {version}"
    info_frame.paragraphs[0].font.name = "Arial"
    info_frame.paragraphs[0].font.size = Pt(18)
    info_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    return slide

def create_section_divider_slide(prs, section_text, page_num):
    """Create DEFRA section divider slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Green background
    bg = slide.shapes.add_shape(1, 0, 0, Cm(25.4), Cm(14.29))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0, 175, 65)
    bg.line.fill.background()

    # White bottom panel
    bottom_panel = slide.shapes.add_shape(1, 0, Cm(11.27), Cm(20.39), Cm(3.08))
    bottom_panel.fill.solid()
    bottom_panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bottom_panel.line.fill.background()

    # White right sidebar
    right_panel = slide.shapes.add_shape(1, Cm(20.39), 0, Cm(5.01), Cm(14.34))
    right_panel.fill.solid()
    right_panel.fill.fore_color.rgb = RGBColor(255, 255, 255)
    right_panel.line.fill.background()

    # DEFRA logo
    add_logo(slide, "assets/logo_tall.png", 21.61, 7.34, 2.29, 3.92)

    # Section header
    header_box = slide.shapes.add_textbox(Cm(1.15), Cm(3.01), Cm(17.75), Cm(1.20))
    header_frame = header_box.text_frame
    header_frame.text = section_text
    header_frame.paragraphs[0].font.name = "Arial"
    header_frame.paragraphs[0].font.size = Pt(22)
    header_frame.paragraphs[0].font.bold = True
    header_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Footer
    footer_box = slide.shapes.add_textbox(Cm(1.27), Cm(13.21), Cm(15.76), Cm(0.76))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Official"
    footer_frame.paragraphs[0].font.name = "Arial"
    footer_frame.paragraphs[0].font.size = Pt(10)
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Page number
    page_box = slide.shapes.add_shape(1, Cm(23.31), Cm(13.24), Cm(1.14), Cm(0.76))
    page_box.fill.solid()
    page_box.fill.fore_color.rgb = RGBColor(0, 176, 80)
    page_box.line.fill.background()

    page_text = page_box.text_frame
    page_text.text = str(page_num)
    page_text.paragraphs[0].font.name = "Arial"
    page_text.paragraphs[0].font.size = Pt(10)
    page_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    page_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def create_content_slide(prs, project_header, slide_title, content_text, page_num):
    """Create DEFRA content slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    bg = slide.shapes.add_shape(1, 0, 0, Cm(25.4), Cm(14.29))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()

    # Green header bar
    header_bar = slide.shapes.add_shape(1, 0, 0, Cm(24.7), Cm(1.05))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(0, 175, 65)
    header_bar.line.fill.background()

    # Header text
    header_box = slide.shapes.add_textbox(Cm(0.18), Cm(0.05), Cm(9.5), Cm(1.03))
    header_frame = header_box.text_frame
    header_frame.text = project_header
    header_frame.paragraphs[0].font.name = "Calibri"
    header_frame.paragraphs[0].font.size = Pt(18)
    header_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Separator line
    sep_line = slide.shapes.add_shape(1, Cm(0.7), Cm(1.10), Cm(24.0), Cm(0.03))
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = RGBColor(0, 175, 65)
    sep_line.line.fill.background()

    # Slide title
    title_box = slide.shapes.add_textbox(Cm(0.70), Cm(1.20), Cm(23.00), Cm(1.30))
    title_frame = title_box.text_frame
    title_frame.text = slide_title
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.size = Pt(18)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Content text
    content_box = slide.shapes.add_textbox(Cm(0.70), Cm(2.60), Cm(23.00), Cm(10.61))
    content_frame = content_box.text_frame
    content_frame.text = content_text
    content_frame.paragraphs[0].font.name = "Calibri"
    content_frame.paragraphs[0].font.size = Pt(14)
    content_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Footer
    footer_box = slide.shapes.add_textbox(Cm(1.27), Cm(13.21), Cm(15.76), Cm(0.76))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Official"
    footer_frame.paragraphs[0].font.name = "Arial"
    footer_frame.paragraphs[0].font.size = Pt(10)
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Page number
    page_box = slide.shapes.add_shape(1, Cm(23.31), Cm(13.24), Cm(1.14), Cm(0.76))
    page_box.fill.solid()
    page_box.fill.fore_color.rgb = RGBColor(0, 176, 80)
    page_box.line.fill.background()

    page_text = page_box.text_frame
    page_text.text = str(page_num)
    page_text.paragraphs[0].font.name = "Arial"
    page_text.paragraphs[0].font.size = Pt(10)
    page_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    page_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def add_placeholder(slide, x, y, w, h, slide_num):
    """Add image placeholder box."""
    placeholder = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(235, 235, 235)
    placeholder.line.color.rgb = RGBColor(0, 175, 65)
    placeholder.line.width = Pt(1.5)

    text_frame = placeholder.text_frame
    text_frame.text = f"[ Insert image — slide {slide_num} ]"
    text_frame.paragraphs[0].font.name = "Arial"
    text_frame.paragraphs[0].font.size = Pt(14)
    text_frame.paragraphs[0].font.color.rgb = RGBColor(102, 102, 102)
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

def create_image_slide(prs, project_header, slide_title, content_text, slide_num, page_num, layout="left"):
    """Create content slide with image placeholder."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    bg = slide.shapes.add_shape(1, 0, 0, Cm(25.4), Cm(14.29))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()

    # Green header bar
    header_bar = slide.shapes.add_shape(1, 0, 0, Cm(24.7), Cm(1.05))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(0, 175, 65)
    header_bar.line.fill.background()

    # Header text
    header_box = slide.shapes.add_textbox(Cm(0.18), Cm(0.05), Cm(9.5), Cm(1.03))
    header_frame = header_box.text_frame
    header_frame.text = project_header
    header_frame.paragraphs[0].font.name = "Calibri"
    header_frame.paragraphs[0].font.size = Pt(18)
    header_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Separator line
    sep_line = slide.shapes.add_shape(1, Cm(0.7), Cm(1.10), Cm(24.0), Cm(0.03))
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = RGBColor(0, 175, 65)
    sep_line.line.fill.background()

    if layout == "left":
        # Image placeholder (left)
        add_placeholder(slide, 0.70, 1.15, 11.80, 12.06, slide_num)

        # Title (right)
        title_box = slide.shapes.add_textbox(Cm(12.90), Cm(1.20), Cm(11.80), Cm(1.30))
        title_frame = title_box.text_frame
        title_frame.text = slide_title
        title_frame.paragraphs[0].font.name = "Arial"
        title_frame.paragraphs[0].font.size = Pt(18)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        # Content (right)
        content_box = slide.shapes.add_textbox(Cm(12.90), Cm(2.60), Cm(11.80), Cm(10.61))
        content_frame = content_box.text_frame
        content_frame.text = content_text
        content_frame.paragraphs[0].font.name = "Calibri"
        content_frame.paragraphs[0].font.size = Pt(14)
        content_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    else:  # layout == "right"
        # Title (left)
        title_box = slide.shapes.add_textbox(Cm(0.70), Cm(1.20), Cm(11.80), Cm(1.30))
        title_frame = title_box.text_frame
        title_frame.text = slide_title
        title_frame.paragraphs[0].font.name = "Arial"
        title_frame.paragraphs[0].font.size = Pt(18)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        # Content (left)
        content_box = slide.shapes.add_textbox(Cm(0.70), Cm(2.60), Cm(11.80), Cm(10.61))
        content_frame = content_box.text_frame
        content_frame.text = content_text
        content_frame.paragraphs[0].font.name = "Calibri"
        content_frame.paragraphs[0].font.size = Pt(14)
        content_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        # Image placeholder (right)
        add_placeholder(slide, 12.90, 1.15, 11.80, 12.06, slide_num)

    # Footer
    footer_box = slide.shapes.add_textbox(Cm(1.27), Cm(13.21), Cm(15.76), Cm(0.76))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Official"
    footer_frame.paragraphs[0].font.name = "Arial"
    footer_frame.paragraphs[0].font.size = Pt(10)
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Page number
    page_box = slide.shapes.add_shape(1, Cm(23.31), Cm(13.24), Cm(1.14), Cm(0.76))
    page_box.fill.solid()
    page_box.fill.fore_color.rgb = RGBColor(0, 176, 80)
    page_box.line.fill.background()

    page_text = page_box.text_frame
    page_text.text = str(page_num)
    page_text.paragraphs[0].font.name = "Arial"
    page_text.paragraphs[0].font.size = Pt(10)
    page_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    page_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def create_two_column_slide(prs, project_header, slide_title, left_content, right_content, page_num):
    """Create two-column content slide."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # White background
    bg = slide.shapes.add_shape(1, 0, 0, Cm(25.4), Cm(14.29))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(255, 255, 255)
    bg.line.fill.background()

    # Green header bar
    header_bar = slide.shapes.add_shape(1, 0, 0, Cm(24.7), Cm(1.05))
    header_bar.fill.solid()
    header_bar.fill.fore_color.rgb = RGBColor(0, 175, 65)
    header_bar.line.fill.background()

    # Header text
    header_box = slide.shapes.add_textbox(Cm(0.18), Cm(0.05), Cm(9.5), Cm(1.03))
    header_frame = header_box.text_frame
    header_frame.text = project_header
    header_frame.paragraphs[0].font.name = "Calibri"
    header_frame.paragraphs[0].font.size = Pt(18)
    header_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    # Separator line
    sep_line = slide.shapes.add_shape(1, Cm(0.7), Cm(1.10), Cm(24.0), Cm(0.03))
    sep_line.fill.solid()
    sep_line.fill.fore_color.rgb = RGBColor(0, 175, 65)
    sep_line.line.fill.background()

    # Full-width title
    title_box = slide.shapes.add_textbox(Cm(0.70), Cm(1.20), Cm(23.00), Cm(1.30))
    title_frame = title_box.text_frame
    title_frame.text = slide_title
    title_frame.paragraphs[0].font.name = "Arial"
    title_frame.paragraphs[0].font.size = Pt(18)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Left column
    left_box = slide.shapes.add_textbox(Cm(0.70), Cm(2.60), Cm(11.30), Cm(10.61))
    left_frame = left_box.text_frame
    left_frame.text = left_content
    left_frame.paragraphs[0].font.name = "Calibri"
    left_frame.paragraphs[0].font.size = Pt(14)
    left_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Right column
    right_box = slide.shapes.add_textbox(Cm(12.40), Cm(2.60), Cm(11.30), Cm(10.61))
    right_frame = right_box.text_frame
    right_frame.text = right_content
    right_frame.paragraphs[0].font.name = "Calibri"
    right_frame.paragraphs[0].font.size = Pt(14)
    right_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Footer
    footer_box = slide.shapes.add_textbox(Cm(1.27), Cm(13.21), Cm(15.76), Cm(0.76))
    footer_frame = footer_box.text_frame
    footer_frame.text = "Official"
    footer_frame.paragraphs[0].font.name = "Arial"
    footer_frame.paragraphs[0].font.size = Pt(10)
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

    # Page number
    page_box = slide.shapes.add_shape(1, Cm(23.31), Cm(13.24), Cm(1.14), Cm(0.76))
    page_box.fill.solid()
    page_box.fill.fore_color.rgb = RGBColor(0, 176, 80)
    page_box.line.fill.background()

    page_text = page_box.text_frame
    page_text.text = str(page_num)
    page_text.paragraphs[0].font.name = "Arial"
    page_text.paragraphs[0].font.size = Pt(10)
    page_text.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    page_text.paragraphs[0].alignment = PP_ALIGN.CENTER

    return slide

def parse_content_input(content_input):
    """Parse user content input into structured slides."""
    slides = []
    lines = content_input.strip().split('\n')

    current_slide = None

    for line in lines:
        line = line.strip()
        if not line:
            continue

        # Check for slide indicators
        if line.startswith('# ') or line.lower().startswith('slide:') or line.lower().startswith('cover:'):
            # New slide
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'type': 'content',
                'title': line.replace('# ', '').replace('Slide:', '').replace('slide:', '').replace('Cover:', '').replace('cover:', '').strip(),
                'content': []
            }
            if line.lower().startswith('cover:'):
                current_slide['type'] = 'cover'

        elif line.startswith('## ') or line.lower().startswith('section:'):
            # Section divider
            if current_slide:
                slides.append(current_slide)
            current_slide = {
                'type': 'section',
                'title': line.replace('## ', '').replace('Section:', '').replace('section:', '').strip(),
                'content': []
            }

        elif line.lower().startswith('[image'):
            # Image requirement
            if current_slide:
                if 'left' in line.lower():
                    current_slide['type'] = 'image_left'
                elif 'right' in line.lower():
                    current_slide['type'] = 'image_right'
                current_slide['content'].append(f"[Image placeholder: {line}]")

        elif line.lower().startswith('[two-column') or line.lower().startswith('[2-column'):
            # Two column layout
            if current_slide:
                current_slide['type'] = 'two_column'

        else:
            # Regular content
            if current_slide:
                current_slide['content'].append(line)
            else:
                # No current slide, create a default one
                current_slide = {
                    'type': 'content',
                    'title': 'Content',
                    'content': [line]
                }

    if current_slide:
        slides.append(current_slide)

    return slides

def create_presentation_from_content(content_input, presentation_title="New DEFRA Presentation", project_name="DEFRA Project"):
    """Create a complete presentation from user content."""
    print(f"Creating DEFRA presentation: {presentation_title}")

    prs = create_defra_presentation()
    slides_data = parse_content_input(content_input)

    # Always start with cover slide
    create_cover_slide(prs, presentation_title)
    print("  Slide 1: Cover slide")

    page_num = 2
    placeholders_report = []
    project_header = f"{project_name} - Report : {datetime.now().strftime('%B %Y')}"

    for slide_idx, slide_data in enumerate(slides_data, 1):
        content_text = '\n'.join(slide_data['content'])

        if slide_data['type'] == 'section':
            create_section_divider_slide(prs, slide_data['title'], page_num)
            print(f"  Slide {page_num}: Section divider - {slide_data['title']}")

        elif slide_data['type'] == 'image_left':
            create_image_slide(prs, project_header, slide_data['title'], content_text, page_num, page_num, "left")
            placeholders_report.append(f"  • Slide {page_num} \"{slide_data['title']}\" — image left")
            print(f"  Slide {page_num}: Content — Image Left - {slide_data['title']}")

        elif slide_data['type'] == 'image_right':
            create_image_slide(prs, project_header, slide_data['title'], content_text, page_num, page_num, "right")
            placeholders_report.append(f"  • Slide {page_num} \"{slide_data['title']}\" — image right")
            print(f"  Slide {page_num}: Content — Image Right - {slide_data['title']}")

        elif slide_data['type'] == 'two_column':
            # Split content into two columns
            content_lines = slide_data['content']
            mid_point = len(content_lines) // 2
            left_content = '\n'.join(content_lines[:mid_point])
            right_content = '\n'.join(content_lines[mid_point:])

            create_two_column_slide(prs, project_header, slide_data['title'], left_content, right_content, page_num)
            print(f"  Slide {page_num}: Two Column Text - {slide_data['title']}")

        else:
            # Standard content slide
            create_content_slide(prs, project_header, slide_data['title'], content_text, page_num)
            print(f"  Slide {page_num}: Content - {slide_data['title']}")

        page_num += 1

    # Save output
    if not os.path.exists("output"):
        os.makedirs("output")

    filename = presentation_title.replace(' ', '_').replace('/', '_') + '.pptx'
    output_path = f"output/{filename}"
    prs.save(output_path)

    print(f"\nPresentation created: {output_path}")

    if placeholders_report:
        print(f"\nℹ Image placeholders added:")
        for placeholder in placeholders_report:
            print(placeholder)

    return output_path

def main():
    """Main creation function."""
    if len(sys.argv) < 2:
        print("Usage: python3 create_defra_deck.py [title] [content_file_or_text]")
        print("\nExample:")
        print("python3 create_defra_deck.py \"My Presentation\" \"content.txt\"")
        return

    title = sys.argv[1]

    if len(sys.argv) > 2:
        content_input = sys.argv[2]
        # Check if it's a file or direct content
        if os.path.exists(content_input):
            with open(content_input, 'r') as f:
                content = f.read()
        else:
            content = content_input
    else:
        # Default sample content
        content = """# Introduction
Welcome to our DEFRA presentation

## Project Overview
This section covers the main project details

# Key Findings
• Finding one
• Finding two
• Finding three

[image left]
# Visual Analysis
This slide needs a chart on the left side
• Supporting point one
• Supporting point two

[two-column]
# Comparison
Left column content:
• Option A benefits
• Option A features

Right column content:
• Option B benefits
• Option B features

# Conclusion
Final thoughts and next steps
"""

    create_presentation_from_content(content, title)

if __name__ == "__main__":
    main()